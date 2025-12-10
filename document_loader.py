import os
from typing import List, Dict, Optional

import docx2txt
from PyPDF2 import PdfReader
from pptx import Presentation

from config import DATA_DIR


class DocumentLoader:
    def __init__(
        self,
        data_dir: str = DATA_DIR,
    ):
        self.data_dir = data_dir
        self.supported_formats = [".pdf", ".pptx", ".docx", ".txt"]

    def load_pdf(self, file_path: str) -> List[Dict]:
        """加载PDF文件，按页返回内容

        TODO: 实现PDF文件加载
        要求：
        1. 使用PdfReader读取PDF文件
        2. 遍历每一页，提取文本内容
        3. 格式化为"--- 第 X 页 ---\n文本内容\n"
        4. 返回pdf内容列表，每个元素包含 {"text": "..."}
        """
        pages = []
        reader = PdfReader(file_path)
        for page_num, page in enumerate(reader.pages, 1):
            text = page.extract_text()
            formatted_text = f"--- 第 {page_num} 页 ---\n{text}\n"
            pages.append({"text": formatted_text})
        return pages

    def load_pptx(self, file_path: str) -> List[Dict]:
        """加载PPT文件，按幻灯片返回内容

        TODO: 实现PPT文件加载
        要求：
        1. 使用Presentation读取PPT文件
        2. 遍历每一页，提取文本内容
        3. 格式化为"--- 幻灯片 X ---\n文本内容\n"
        4. 返回幻灯片内容列表，每个元素包含 {"text": "..."}
        """
        slides = []
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text)
            text = "\n".join(text_parts)
            formatted_text = f"--- 幻灯片 {slide_num} ---\n{text}\n"
            slides.append({"text": formatted_text})
        return slides

    def load_docx(self, file_path: str) -> str:
        """加载DOCX文件
        TODO: 实现DOCX文件加载
        要求：
        1. 使用docx2txt读取DOCX文件
        2. 返回文本内容
        """
        text = docx2txt.process(file_path)
        return text

    def load_txt(self, file_path: str) -> str:
        """加载TXT文件
        TODO: 实现TXT文件加载
        要求：
        1. 使用open读取TXT文件（注意使用encoding="utf-8"）
        2. 返回文本内容
        """
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
        return text

    def load_document(self, file_path: str) -> List[Dict[str, str]]:
        """加载单个文档，PDF和PPT按页/幻灯片分割，返回文档块列表"""
        ext = os.path.splitext(file_path)[1].lower()
        filename = os.path.basename(file_path)
        documents = []

        if ext == ".pdf":
            pages = self.load_pdf(file_path)
            for page_idx, page_data in enumerate(pages, 1):
                documents.append(
                    {
                        "content": page_data["text"],
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": page_idx,
                    }
                )
        elif ext == ".pptx":
            slides = self.load_pptx(file_path)
            for slide_idx, slide_data in enumerate(slides, 1):
                documents.append(
                    {
                        "content": slide_data["text"],
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": slide_idx,
                    }
                )
        elif ext == ".docx":
            content = self.load_docx(file_path)
            if content:
                documents.append(
                    {
                        "content": content,
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": 0,
                    }
                )
        elif ext == ".txt":
            content = self.load_txt(file_path)
            if content:
                documents.append(
                    {
                        "content": content,
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": 0,
                    }
                )
        else:
            print(f"不支持的文件格式: {ext}")

        return documents

    def load_all_documents(self) -> List[Dict[str, str]]:
        """加载数据目录下的所有文档"""
        if not os.path.exists(self.data_dir):
            print(f"数据目录不存在: {self.data_dir}")
            return None

        documents = []

        for root, dirs, files in os.walk(self.data_dir):
            for file in files:
                ext = os.path.splitext(file)[1].lower()
                if ext in self.supported_formats:
                    file_path = os.path.join(root, file)
                    print(f"正在加载: {file_path}")
                    doc_chunks = self.load_document(file_path)
                    if doc_chunks:
                        documents.extend(doc_chunks)

        return documents
